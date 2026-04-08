"""
高还原度 PPTX 生成器

策略：
  1. Gemini 视觉模型分析原始信息图，提取所有文本元素的位置和样式
  2. Gemini 图片编辑模型擦除原图上的文字，生成干净背景图
  3. 干净背景图作为幻灯片全屏背景
  4. 在文字原位叠加可编辑的透明文本框

用法:
    # 单张图片:
    python generate_pptx.py --image <图片路径> --output output.pptx

    # 批量处理（目录下所有图片按文件名排序）:
    python generate_pptx.py --image-dir <图片目录> --output output.pptx

    # 示例:
    python generate_pptx.py --image-dir examples/俄罗斯地理和文化 --output output.pptx --keep-clean
"""

import argparse
import os
import re
import sys
import json
import tempfile
from pathlib import Path

from src.gemini_client import generate_text_with_images, edit_image

DEFAULT_MODEL = "gemini-3.1-pro-preview"
IMAGE_EDIT_MODEL = "gemini-3.1-flash-image-preview"

# ═══════════════════════════════════════════════════════════
#  Prompt：让 LLM 分析图片中的文本元素并返回结构化数据
# ═══════════════════════════════════════════════════════════

SYSTEM_PROMPT = """\
你是一位专业的幻灯片设计分析师。你的任务是分析信息图图片，提取其中所有文本元素的精确位置、样式信息。

## 输出格式

返回一个 JSON 对象，包含以下字段：

```json
{
  "text_elements": [
    {
      "text": "文本内容",
      "left_pct": 15.0,
      "top_pct": 8.0,
      "width_pct": 70.0,
      "height_pct": 10.0,
      "font_size_pt": 40,
      "font_bold": true,
      "font_color_hex": "FFFFFF",
      "alignment": "center"
    }
  ]
}
```

## 字段说明

- `left_pct`, `top_pct`: 文本框左上角相对于图片宽高的百分比（0-100）。
- `width_pct`, `height_pct`: 文本框宽高相对于图片宽高的百分比（0-100）。
  这四个值要精确地框住文字区域（紧贴文字边界，上下左右各留约 0.5% 余量即可）。
- `font_size_pt`: 估算的字号（磅值），按 13.333×7.5 英寸幻灯片估算
- `font_bold`: 是否粗体
- `font_color_hex`: 字体颜色（6位十六进制，不含#）
- `alignment`: "left" | "center" | "right"

## 额外字段

在 JSON 根对象中还需包含：

- `graphic_elements`: 字符串数组，列出图片中所有非文字的视觉元素（如 "俄罗斯地图轮廓", "红白蓝三色飘带", "地球图标", "桥梁图标", "时钟图标", "渐变背景", "水平装饰线" 等）

## 要求

1. 提取图片中 **所有** 可见的文本元素，包括标题、副标题、正文、数据数字、脚注、页码等
2. bounding box 要尽可能紧贴文字
3. 颜色要尽可能准确
4. `graphic_elements` 要列出所有重要的非文字视觉元素
5. 只输出 JSON，不要输出任何解释
6. JSON 用 ```json ``` 包裹
"""

USER_PROMPT = "请分析这张信息图幻灯片中的所有文本元素，提取位置、样式信息。返回 JSON 格式数据。"

TEXT_REMOVAL_PROMPT_TEMPLATE = (
    "Edit this image to remove ALL text and numbers. "
    "Here is the complete list of text to erase:\n"
    "{text_list}\n\n"
    "IMPORTANT:\n"
    "- Erase every single item listed above, including large styled numbers like '1709万', '9000', '14' etc.\n"
    "- Numbers are TEXT, not graphics — they MUST be removed\n"
    "- After removal, fill each area with the surrounding background color/texture seamlessly\n"
    "- KEEP these graphic elements untouched: {keep_elements}\n"
    "- Do NOT alter colors, layout, or any non-text visual elements\n"
    "- Output at the same resolution as input"
)


# ═══════════════════════════════════════════════════════════
#  JSON 解析
# ═══════════════════════════════════════════════════════════

def _extract_json(response_text: str) -> dict:
    """从 LLM 回复中提取 JSON 数据。"""
    match = re.search(r'```json\s*\n(.*?)```', response_text, re.DOTALL)
    json_str = match.group(1).strip() if match else response_text.strip()
    return json.loads(json_str)


# ═══════════════════════════════════════════════════════════
#  Gemini 图片文字擦除
# ═══════════════════════════════════════════════════════════

VERIFY_PROMPT = '这张图片中是否还有任何可见的文字、数字或字符？如果有，请逐条列出剩余的文字内容。如果没有任何文字，只回复"无"。'


def _verify_text_removed(image_path: str, model: str = DEFAULT_MODEL) -> list[str]:
    """检查图片中是否还有残留文字，返回残留文字列表。"""
    response = generate_text_with_images(
        model=model,
        system_prompt="你是一个图片文字检测工具。仔细检查图片中是否还有任何文字或数字。",
        user_prompt=VERIFY_PROMPT,
        image_paths=[image_path],
    )
    text = response.strip()
    if text == "无" or "没有" in text or "无任何" in text or "no text" in text.lower():
        return []
    # 提取残留文字列表
    remaining = []
    for line in text.split("\n"):
        line = line.strip().lstrip("*-•· 0123456789.)")
        if line:
            remaining.append(line)
    return remaining


def remove_text_from_image(
    image_path: str,
    output_path: str,
    text_elements: list[dict],
    graphic_elements: list[str] | None = None,
    model: str = IMAGE_EDIT_MODEL,
    max_retries: int = 2,
) -> str:
    """
    使用 Gemini 图片编辑模型擦除图片上的所有文字。
    擦除后会验证是否有残留文字，如有则重试。

    Args:
        image_path: 原始图片路径
        output_path: 输出干净图片的路径
        text_elements: 需要擦除的文本元素列表
        graphic_elements: 需要保留的图形元素描述列表
        model: 图片编辑模型名称
        max_retries: 最大重试次数

    Returns:
        输出图片的路径
    """
    # 构建要保留的图形元素描述
    if graphic_elements:
        keep_elements = ", ".join(graphic_elements)
    else:
        keep_elements = "background, icons, map, lines, ribbons, all non-text graphics"

    current_input = image_path
    current_texts = text_elements

    for attempt in range(1, max_retries + 2):
        # 构建要擦除的文本列表
        text_list = "\n".join(
            f'- "{elem["text"] if isinstance(elem, dict) else elem}"'
            for elem in current_texts
        )

        prompt = TEXT_REMOVAL_PROMPT_TEMPLATE.format(
            text_list=text_list,
            keep_elements=keep_elements,
        )

        print(f"  [Gemini] 正在擦除文字（第 {attempt} 轮）...")
        image_data = edit_image(
            model=model,
            prompt=prompt,
            image_path=current_input,
        )
        if image_data is None:
            if attempt == 1:
                raise RuntimeError(
                    "Gemini 图片编辑未返回图片数据。\n"
                    "可能原因：\n"
                    "1. 图片编辑模型名称不正确（当前: {}）\n"
                    "2. 第三方代理不支持图片编辑功能\n"
                    "3. API 权限不足\n\n"
                    "建议：尝试使用 Google 官方 API，或启用「跳过文本叠加层」选项".format(model)
                )
            else:
                print(f"  ⚠ 第 {attempt} 轮编辑失败，使用上一次结果")
                break

        with open(output_path, "wb") as f:
            f.write(image_data)

        # 验证是否还有残留文字
        if attempt <= max_retries:
            print(f"  [验证] 检查残留文字...")
            remaining = _verify_text_removed(output_path)
            if not remaining:
                print(f"  ✓ 文字已全部擦除")
                break
            else:
                print(f"  ✗ 仍有残留文字: {remaining}")
                current_input = output_path
                current_texts = remaining
        else:
            print(f"  已达最大重试次数")

    return output_path


# ═══════════════════════════════════════════════════════════
#  LLM 图片分析
# ═══════════════════════════════════════════════════════════

def analyze_image_text(image_path: str, model: str = DEFAULT_MODEL) -> dict:
    """使用 LLM 分析图片中的文本元素。"""
    print(f"  [LLM] 正在分析图片文本: {image_path}")
    response = generate_text_with_images(
        model=model,
        system_prompt=SYSTEM_PROMPT,
        user_prompt=USER_PROMPT,
        image_paths=[image_path],
    )
    return _extract_json(response)


# ═══════════════════════════════════════════════════════════
#  核心：生成 PPTX
# ═══════════════════════════════════════════════════════════

def build_slide_from_image(
    prs,
    image_path: str,
    text_elements: list[dict] | None = None,
    slide_width_inches: float = 13.333,
    slide_height_inches: float = 7.5,
):
    """
    创建一张以（干净）图片为背景的幻灯片，并叠加可编辑文本框。
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    sw = Inches(slide_width_inches)
    sh = Inches(slide_height_inches)

    # ── 1. 添加全屏背景图片（已擦除文字） ──
    slide.shapes.add_picture(image_path, 0, 0, sw, sh)

    # ── 2. 叠加可编辑文本框 ──
    if text_elements:
        align_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }

        for elem in text_elements:
            left = Emu(int(sw * elem["left_pct"] / 100))
            top = Emu(int(sh * elem["top_pct"] / 100))
            width = Emu(int(sw * elem["width_pct"] / 100))
            height = Emu(int(sh * elem["height_pct"] / 100))

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0

            p = tf.paragraphs[0]
            p.text = elem["text"]
            p.alignment = align_map.get(elem.get("alignment", "left"), PP_ALIGN.LEFT)

            font = p.font
            font.size = Pt(int(elem.get("font_size_pt", 14) * 0.85))
            font.bold = elem.get("font_bold", False)
            font.name = "Microsoft YaHei"

            color_hex = elem.get("font_color_hex", "333333")
            font.color.rgb = RGBColor(
                int(color_hex[0:2], 16),
                int(color_hex[2:4], 16),
                int(color_hex[4:6], 16),
            )

    return slide


def process_single_image(
    prs,
    image_path: str,
    model: str = DEFAULT_MODEL,
    image_edit_model: str = IMAGE_EDIT_MODEL,
    skip_text_overlay: bool = False,
    keep_clean_image: bool = False,
):
    """处理单张图片：分析文本 → Gemini 擦除文字 → 构建幻灯片。"""
    text_elements = None
    clean_image_path = image_path  # 默认使用原图

    if not skip_text_overlay:
        try:
            # Step 1: 分析文本元素
            data = analyze_image_text(image_path, model)
            text_elements = data.get("text_elements", [])
            print(f"  提取到 {len(text_elements)} 个文本元素")

            # Step 2: 用 Gemini 擦除文字
            graphic_elements = data.get("graphic_elements", None)
            if text_elements:
                if keep_clean_image:
                    base, ext = os.path.splitext(image_path)
                    clean_image_path = f"{base}_clean.png"
                else:
                    fd, clean_image_path = tempfile.mkstemp(suffix=".png")
                    os.close(fd)

                remove_text_from_image(
                    image_path, clean_image_path,
                    text_elements, graphic_elements, image_edit_model,
                )
                print(f"  干净背景图已生成: {clean_image_path}")
        except Exception as e:
            print(f"  [警告] 处理失败，将使用原图: {e}")
            import traceback
            traceback.print_exc()
            clean_image_path = image_path

    # Step 3: 构建幻灯片
    build_slide_from_image(prs, clean_image_path, text_elements)

    # 清理临时文件
    if clean_image_path != image_path and not keep_clean_image:
        try:
            os.unlink(clean_image_path)
        except OSError:
            pass


def generate_pptx_from_image(
    image_path: str,
    output_path: str,
    model: str = DEFAULT_MODEL,
    image_edit_model: str = IMAGE_EDIT_MODEL,
    skip_text_overlay: bool = False,
    keep_clean_image: bool = False,
) -> bool:
    """从单张图片生成 PPTX。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    process_single_image(
        prs, image_path, model, image_edit_model,
        skip_text_overlay, keep_clean_image,
    )

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    print(f"PPTX 已保存: {output_path}")
    return True


def generate_pptx_from_images(
    image_paths: list[str],
    output_path: str,
    model: str = DEFAULT_MODEL,
    image_edit_model: str = IMAGE_EDIT_MODEL,
    skip_text_overlay: bool = False,
    keep_clean_image: bool = False,
) -> bool:
    """从多张图片生成包含多页的 PPTX。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, img_path in enumerate(image_paths):
        print(f"\n── 处理第 {i + 1}/{len(image_paths)} 张 ──")
        process_single_image(
            prs, img_path, model, image_edit_model,
            skip_text_overlay, keep_clean_image,
        )

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    print(f"\nPPTX 已保存: {output_path} ({len(image_paths)} 页)")
    return True


# ═══════════════════════════════════════════════════════════
#  主入口
# ═══════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(
        description="高还原度 PPTX 生成（Gemini 擦除文字 + 图片背景 + 文本叠加）"
    )
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--image", help="单张信息图图片路径")
    group.add_argument("--image-dir", help="信息图图片目录（按文件名排序）")
    parser.add_argument("--output", default="output.pptx", help="输出 PPTX 路径")
    parser.add_argument("--model", default=DEFAULT_MODEL, help="文本分析模型")
    parser.add_argument(
        "--image-edit-model", default=IMAGE_EDIT_MODEL,
        help="图片编辑模型（用于擦除文字）"
    )
    parser.add_argument(
        "--no-text", action="store_true",
        help="不添加文本叠加层（仅图片背景）"
    )
    parser.add_argument(
        "--keep-clean", action="store_true",
        help="保留擦除文字后的干净图片（默认删除临时文件）"
    )

    args = parser.parse_args()

    if args.image:
        success = generate_pptx_from_image(
            args.image, args.output, args.model,
            args.image_edit_model, args.no_text, args.keep_clean,
        )
    else:
        img_dir = Path(args.image_dir)
        exts = {".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff"}
        images = sorted(
            [str(p) for p in img_dir.iterdir() if p.suffix.lower() in exts]
        )
        if not images:
            print(f"目录中未找到图片: {args.image_dir}")
            sys.exit(1)
        print(f"找到 {len(images)} 张图片")
        success = generate_pptx_from_images(
            images, args.output, args.model,
            args.image_edit_model, args.no_text, args.keep_clean,
        )

    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
